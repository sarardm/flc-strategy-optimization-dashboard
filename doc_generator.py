"""
Document Generator for FLC Portfolio Optimization Dashboard
============================================================
Generates executive summary .docx and slide deck .pptx files
for each Phase 1 framework analysis.
"""

import os
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.dml.color import RGBColor as PptxRGB
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu

from data import (
    INSTITUTION, PESTLE_DATA, PORTERS_DATA, PORTERS_INSIGHTS,
    BCG_DATA, BCG_DEPT_DATA, BCG_DEPT_INSIGHTS, BCG_INSIGHTS, BCG_QUADRANT_COLORS,
    GRAY_ASSOCIATES_DATA, GA_INSIGHTS, FRAMEWORK_DESCRIPTIONS,
    SWOT_DATA, ZONE_TO_WIN_DATA, SCENARIOS, RISK_MITIGATION,
    ROADMAP_MILESTONES, ROADMAP_KPIS,
)

OUTPUT_DIR = os.path.join(os.path.dirname(__file__), "generated_docs")
os.makedirs(OUTPUT_DIR, exist_ok=True)

FLC_NAVY = RGBColor(0x00, 0x30, 0x57)
FLC_GOLD = RGBColor(0xC8, 0xA4, 0x15)


# ============================================================================
# SHARED HELPERS
# ============================================================================

def _add_heading(doc, text, level=1):
    h = doc.add_heading(text, level=level)
    for run in h.runs:
        run.font.color.rgb = FLC_NAVY
    return h


def _add_para(doc, text, bold=False, size=11):
    p = doc.add_paragraph()
    run = p.add_run(text)
    run.font.size = Pt(size)
    run.bold = bold
    return p


def _doc_header(doc, title, subtitle):
    doc.add_paragraph()
    h = doc.add_heading(title, level=0)
    for run in h.runs:
        run.font.color.rgb = FLC_NAVY
    p = doc.add_paragraph()
    run = p.add_run(subtitle)
    run.font.size = Pt(12)
    run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    p = doc.add_paragraph()
    run = p.add_run(f"Fort Lewis College | {INSTITUTION['location']} | {INSTITUTION['type']}")
    run.font.size = Pt(10)
    run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    doc.add_paragraph("_" * 60)
    return doc


def _pptx_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide


def _pptx_content_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        if i == 0:
            tf.text = bullet
        else:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
    return slide


# ============================================================================
# PESTLE DOCUMENTS
# ============================================================================

def generate_pestle_docx():
    doc = Document()
    _doc_header(doc, "PESTLE Analysis: Executive Summary",
                "External Environmental Scan for Fort Lewis College Academic Affairs")

    _add_heading(doc, "1. Introduction & Methodology", 2)
    _add_para(doc, FRAMEWORK_DESCRIPTIONS["PESTLE"])
    _add_para(doc, (
        "This analysis draws on FLC's internal PESTLE report, the External Forces Shaping "
        "Fort Lewis College presentation, institutional enrollment data, and the 2025-2030 "
        "Strategic Plan. Each of the six PESTLE dimensions was assessed for impact severity "
        "(1-5 scale) and directional trend (Positive, Mixed, Negative, Stable, Opportunity)."
    ))

    _add_heading(doc, "2. Key Findings by Category", 2)
    for category, d in PESTLE_DATA.items():
        _add_heading(doc, f"{category} (Impact: {d['impact']}, Trend: {d['trend']})", 3)
        for factor in d["factors"]:
            doc.add_paragraph(factor, style="List Bullet")
        _add_para(doc, "Opportunities:", bold=True, size=10)
        for opp in d["opportunities"]:
            doc.add_paragraph(opp, style="List Bullet 2")

    _add_heading(doc, "3. Impact Assessment Summary", 2)
    _add_para(doc, (
        "The highest-impact categories are Economic (5/5) and Political/Social (4/5 each). "
        "Economic factors\u2014particularly state funding volatility, tuition sensitivity, and the "
        "revenue impact of the Native American tuition waiver\u2014represent the most significant "
        "external pressures on FLC's financial model. Political factors including performance-based "
        "funding and pressure on DEI programs add further uncertainty. Social factors, notably "
        "declining college-going rates and shifting student expectations toward career outcomes, "
        "require programmatic adaptation."
    ))
    _add_para(doc, (
        "Technological factors (4/5 impact) reflect a rapidly changing landscape: the online market is "
        "saturated and FLC has no online brand, though AI is transforming delivery. Legal factors (4/5) "
        "are deteriorating with NATW legal basis (CRS 23-52-105) requiring proactive documentation. "
        "Environmental factors present both risk (wildfire/drought) and opportunity (outdoor brand)."
    ))

    _add_heading(doc, "4. Strategic Implications for Academic Affairs", 2)
    _add_para(doc, (
        "The PESTLE analysis reveals that FLC operates in an environment of heightened political risk "
        "and constrained resources. Federal DEI policy disruptions threaten TRIO programs and could "
        "misclassify the statutory NATW mission. The Durango housing crisis constrains faculty/staff "
        "recruitment. Key strategic imperatives include: (1) reframing Indigenous education through "
        "statutory obligations (CRS 23-52-105) and state law (not DEI), (2) strengthening the dual enrollment pipeline "
        "as a hedge against declining first-year enrollment, (3) improving retention as the most "
        "cost-effective enrollment strategy, and (4) investing in AI capabilities while recognizing "
        "that large-scale online expansion faces a saturated market where FLC has no brand."
    ))

    _add_heading(doc, "5. Recommendations", 2)
    recommendations = [
        "Proactively document NATW statutory basis (CRS 23-52-105) to protect against DEI misclassification",
        "Prioritize retention improvement as the most cost-effective enrollment strategy (Compass, early-alert)",
        "Grow dual enrollment pipeline and transfer pathways as near-term enrollment stabilizers",
        "Frame Indigenous education through statutory obligations (CRS 23-52-105) and cultural preservation, not DEI language",
        "Invest in AI Institute and experiential learning as place-based institutional differentiators",
        "Qualify online expansion: pursue only Indigenous niche (NATW moat), not generic online degrees",
        "Address Durango housing crisis impact on faculty/staff recruitment through institutional partnerships",
    ]
    for rec in recommendations:
        doc.add_paragraph(rec, style="List Bullet")

    _add_para(doc, "\nSource: FLC PESTLE Report, External Forces Shaping FLC presentation, "
              "FLC Institutional Data, 2025-2030 Strategic Plan", size=9)

    path = os.path.join(OUTPUT_DIR, "PESTLE_Executive_Summary.docx")
    doc.save(path)
    return path


def generate_pestle_pptx():
    prs = Presentation()
    _pptx_title_slide(prs, "PESTLE Analysis", "External Environmental Scan\nFort Lewis College Academic Affairs")

    _pptx_content_slide(prs, "Methodology", [
        FRAMEWORK_DESCRIPTIONS["PESTLE"],
        "Sources: PESTLE_Report_FLC.docx, External Forces Shaping FLC.pptx, Institutional Data",
        "Six dimensions assessed on 1-5 impact scale with directional trends",
    ])

    for category, d in PESTLE_DATA.items():
        bullets = [f"Impact: {d['impact']} ({d['impact_score']}/5) | Trend: {d['trend']}"]
        bullets.extend(d["factors"][:4])
        bullets.append(f"Opportunity: {d['opportunities'][0]}")
        _pptx_content_slide(prs, f"{category} Factors", bullets)

    _pptx_content_slide(prs, "Impact Assessment Summary", [
        "Highest impact: Political (5/5), Economic (5/5)",
        "High impact: Social (4/5), Technological (4/5), Legal (4/5)",
        "Key risks: Federal DEI policy, tribal waiver vulnerability, Durango housing crisis",
        "Key opportunity: Indigenous education (statutorily grounded), AI Institute, experiential learning",
    ])

    _pptx_content_slide(prs, "Strategic Recommendations", [
        "Protect NATW statutory basis (CRS 23-52-105) from DEI misclassification",
        "Prioritize retention as most cost-effective enrollment strategy",
        "Frame Indigenous education through statutory/sovereign obligations, not DEI language",
        "Invest in AI Institute and place-based experiential learning",
        "Qualify online expansion: Indigenous niche only (NATW moat), not generic degrees",
    ])

    path = os.path.join(OUTPUT_DIR, "PESTLE_Slide_Deck.pptx")
    prs.save(path)
    return path


# ============================================================================
# PORTER'S FIVE FORCES DOCUMENTS
# ============================================================================

def generate_porters_docx():
    doc = Document()
    _doc_header(doc, "Porter's Five Forces: Executive Summary",
                "Competitive Analysis of FLC's Higher Education Market")

    _add_heading(doc, "1. Introduction & Methodology", 2)
    _add_para(doc, FRAMEWORK_DESCRIPTIONS["Porters"])
    _add_para(doc, (
        "This analysis applies Porter's framework using published higher education competitive "
        "research methodologies combined with FLC institutional data, enrollment trends, and "
        "regional market intelligence. Each force is rated on a 1-5 scale (Low to High competitive "
        "intensity) with supporting indicators and trend analysis."
    ))

    _add_heading(doc, "2. Analysis of Five Forces", 2)
    for force, d in PORTERS_DATA.items():
        _add_heading(doc, f"{force}: {d['rating']} ({d['score']}/5)", 3)
        _add_para(doc, d["description"])
        _add_para(doc, "Key Indicators:", bold=True, size=10)
        for ind in d["indicators"]:
            doc.add_paragraph(
                f"{ind['name']}: {ind['value']} (Trend: {ind['trend']})",
                style="List Bullet"
            )

    _add_heading(doc, "3. Overall Competitive Assessment", 2)
    _add_para(doc, (
        "The aggregate competitive intensity facing Fort Lewis College is HIGH, with an average "
        "force score of 3.4/5 across all five dimensions. Competitive Rivalry (4.0/5) and Bargaining "
        "Power of Students (4.0/5) represent the most intense pressures. Online competition, while "
        "significant nationally, is unverified specifically for FLC's student population, which is "
        "predominantly place-bound and values the residential outdoor experience."
    ))
    _add_para(doc, (
        "FLC's strongest defensive positions are its unique Native American mission and tuition "
        "waiver program (which creates a competitive moat for serving Indigenous students), the "
        "Durango outdoor recreation lifestyle (which cannot be replicated by online competitors), "
        "and the intimate small-college experience with 19-student average class sizes."
    ))

    _add_heading(doc, "4. Strategic Implications", 2)
    for insight in PORTERS_INSIGHTS:
        doc.add_paragraph(insight, style="List Bullet")

    _add_heading(doc, "5. Competitive Positioning Recommendations", 2)
    recommendations = [
        "Double down on experiential/outdoor differentiators that online competitors cannot replicate",
        "Develop strategic pricing and financial aid models to address student price sensitivity",
        "Build transfer-friendly pathways with community colleges rather than competing against them",
        "Invest in unique program offerings (Indigenous education, sustainability, AI) with limited regional competition",
        "Address faculty recruitment challenges through creative compensation and quality-of-life incentives",
        "Expand online offerings strategically to meet student flexibility expectations while preserving residential core",
    ]
    for rec in recommendations:
        doc.add_paragraph(rec, style="List Bullet")

    _add_para(doc, "\nSource: Porter's Five Forces framework (published research); "
              "Applied to FLC institutional data and regional market intelligence", size=9)

    path = os.path.join(OUTPUT_DIR, "Porters_Executive_Summary.docx")
    doc.save(path)
    return path


def generate_porters_pptx():
    prs = Presentation()
    _pptx_title_slide(prs, "Porter's Five Forces Analysis",
                       "Competitive Market Assessment\nFort Lewis College Academic Affairs")

    _pptx_content_slide(prs, "Methodology", [
        FRAMEWORK_DESCRIPTIONS["Porters"],
        "Source: Published higher education research + FLC institutional data",
        "Each force rated 1-5 (Low-High intensity) with indicators and trends",
    ])

    for force, d in PORTERS_DATA.items():
        bullets = [f"Rating: {d['rating']} ({d['score']}/5)"]
        bullets.append(d["description"])
        for ind in d["indicators"][:3]:
            bullets.append(f"{ind['name']}: {ind['value']} ({ind['trend']})")
        _pptx_content_slide(prs, force, bullets)

    _pptx_content_slide(prs, "Overall Assessment", [
        "Average competitive intensity: 3.8/5 (HIGH)",
        "Strongest pressures: Competitive Rivalry (4.5), Student Power (4.0)",
        "FLC defensive moats: Native American mission, outdoor lifestyle, small-college experience",
        "Greatest vulnerabilities: online competition, price sensitivity, faculty recruitment",
    ])

    _pptx_content_slide(prs, "Strategic Recommendations", [
        "Leverage experiential/outdoor differentiators against online competitors",
        "Develop strategic pricing models for student price sensitivity",
        "Build transfer pathways with community colleges as partners, not rivals",
        "Invest in unique programs with limited regional competition",
        "Address faculty recruitment through creative compensation",
    ])

    path = os.path.join(OUTPUT_DIR, "Porters_Slide_Deck.pptx")
    prs.save(path)
    return path


# ============================================================================
# GRAY ASSOCIATES DOCUMENTS
# ============================================================================

def generate_gray_docx():
    doc = Document()
    _doc_header(doc, "Gray Associates Portfolio Analysis: Executive Summary",
                "Academic Program Evaluation for Fort Lewis College")

    _add_heading(doc, "1. Introduction & Methodology", 2)
    _add_para(doc, FRAMEWORK_DESCRIPTIONS["Gray"])
    _add_para(doc, (
        "This analysis applies the Gray Associates Program Evaluation System (PES) methodology "
        "to FLC's academic programs using institutional enrollment data, BCG market share analysis, "
        "and regional employment projections. Market Score is a weighted composite of Student "
        "Demand (40%), Employment Outlook (40%), and Competitive Position (20%). Economics Score "
        "is derived from SCH generation efficiency and program cost structures."
    ))
    _add_para(doc, (
        "Note: This analysis uses the Gray Associates framework methodology sourced from published "
        "research and applied to FLC's internal institutional data. FLC does not have a Gray Associates "
        "subscription; scores represent best estimates based on available data."
    ), size=10)

    _add_heading(doc, "2. Program Portfolio Classification", 2)
    categories = {"Grow": [], "Sustain": [], "Transform": [], "Evaluate": [], "Sunset Review": []}
    for _, row in GRAY_ASSOCIATES_DATA.iterrows():
        categories[row["GA_Recommendation"]].append(row)

    for cat, programs in categories.items():
        if programs:
            _add_heading(doc, f"{cat} ({len(programs)} programs)", 3)
            for p in programs:
                doc.add_paragraph(
                    f"{p['Program']}: Market Score {p['Market_Score']}, Economics {p['Economics_Score']}, "
                    f"Enrollment {p['Enrollment']}, Mission: {p['Mission_Alignment']}",
                    style="List Bullet"
                )

    _add_heading(doc, "3. Key Findings", 2)
    for insight in GA_INSIGHTS:
        doc.add_paragraph(insight, style="List Bullet")

    _add_heading(doc, "4. Market Demand Analysis", 2)
    _add_para(doc, (
        "Programs with the highest Market Scores reflect strong alignment between student interest "
        "and labor market demand. Engineering (79), Health Sciences (76), Computer Information "
        "Systems (75), and Business Administration (74) lead on market positioning. These programs "
        "benefit from robust regional and national employment demand in STEM, healthcare, and "
        "business fields."
    ))
    _add_para(doc, (
        "Programs with the lowest Market Scores\u2014Philosophy (34), Political Science (38), "
        "Anthropology (38), Art & Design (39), and History (39)\u2014face challenging demand dynamics. "
        "However, some of these programs play important roles in general education requirements and "
        "institutional mission fulfillment that pure market metrics do not capture."
    ))

    _add_heading(doc, "5. Recommendations", 2)
    recommendations = [
        "Invest aggressively in Grow programs: add capacity, online options, and graduate tracks",
        "Improve efficiency of Sustain programs through shared resources and cross-listing",
        "Innovate delivery for Transform programs: English and Math could adopt co-requisite and applied models",
        "Conduct deep-dive reviews of Evaluate and Sunset Review programs with faculty input",
        "Consider interdisciplinary consolidation for low-enrollment programs to preserve breadth while reducing cost",
        "Weight mission alignment alongside market/economics in final program decisions",
    ]
    for rec in recommendations:
        doc.add_paragraph(rec, style="List Bullet")

    _add_para(doc, "\nSource: Gray Associates PES methodology (internet research); "
              "Applied to FLC enrollment data, BCG analysis, and regional employment data", size=9)

    path = os.path.join(OUTPUT_DIR, "Gray_Executive_Summary.docx")
    doc.save(path)
    return path


def generate_gray_pptx():
    prs = Presentation()
    _pptx_title_slide(prs, "Gray Associates Portfolio Analysis",
                       "Academic Program Evaluation\nFort Lewis College")

    _pptx_content_slide(prs, "Methodology", [
        FRAMEWORK_DESCRIPTIONS["Gray"],
        "Market Score = Student Demand (40%) + Employment (40%) + Competition (20%)",
        "Economics Score = SCH generation efficiency + program cost structure",
        "Data Source: Gray Associates PES methodology applied to FLC institutional data",
    ])

    for cat in ["Grow", "Sustain", "Transform", "Evaluate", "Sunset Review"]:
        progs = GRAY_ASSOCIATES_DATA[GRAY_ASSOCIATES_DATA["GA_Recommendation"] == cat]
        if not progs.empty:
            bullets = [f"{len(progs)} programs classified as {cat}"]
            for _, p in progs.iterrows():
                bullets.append(f"{p['Program']} (Market: {p['Market_Score']}, Econ: {p['Economics_Score']})")
            _pptx_content_slide(prs, f"{cat} Programs", bullets[:6])

    _pptx_content_slide(prs, "Key Findings", GA_INSIGHTS)

    _pptx_content_slide(prs, "Recommendations", [
        "Invest in Grow programs: capacity, online, graduate tracks",
        "Improve efficiency of Sustain programs via shared resources",
        "Innovate delivery for Transform programs (English, Math)",
        "Deep-dive reviews for Evaluate/Sunset programs with faculty",
        "Weight mission alignment in final decisions",
    ])

    path = os.path.join(OUTPUT_DIR, "Gray_Slide_Deck.pptx")
    prs.save(path)
    return path


# ============================================================================
# BCG MATRIX DOCUMENTS
# ============================================================================

def generate_bcg_docx():
    doc = Document()
    _doc_header(doc, "BCG Growth-Share Matrix: Executive Summary",
                "Department & Major-Level Portfolio Analysis for Fort Lewis College")

    _add_heading(doc, "1. Introduction & Methodology", 2)
    _add_para(doc, FRAMEWORK_DESCRIPTIONS["BCG"])
    _add_para(doc, (
        "This analysis examines FLC's portfolio at two levels. The department-level view (22 departments) "
        "uses % of total Student Credit Hours (SCH) as the market share proxy, providing a revenue-oriented "
        "perspective. The major-level view (48 majors) uses 2024 enrollment headcount vs. 2022\u20132024 "
        "percentage change, with bubble size reflecting absolute students gained or lost."
    ))

    # ── Department-Level Analysis ──
    _add_heading(doc, "2. Department-Level Analysis (SCH-Based)", 2)
    _add_para(doc, (
        "The department-level BCG matrix plots 22 departments by their share of total Student Credit "
        "Hours (X-axis) against 2-year enrollment change (Y-axis). The SCH divider is 4.0%, separating "
        "departments that generate a meaningful share of institutional teaching revenue from those that do not."
    ))
    dept_quadrants = {"Star": "Stars", "Cash Cow": "Cash Cows",
                      "Question Mark": "Question Marks", "Concern": "Concerns"}
    for q_key, q_label in dept_quadrants.items():
        depts = BCG_DEPT_DATA[BCG_DEPT_DATA["Quadrant"] == q_key]
        _add_heading(doc, f"{q_label} ({len(depts)} departments)", 3)
        for _, row in depts.iterrows():
            doc.add_paragraph(
                f"{row['Department']}: {row['SCH_Pct']}% of SCH, {row['Two_Year_Change']:+.1f}% 2-year change",
                style="List Bullet"
            )
    _add_para(doc, "Department-Level Insights:", bold=True, size=10)
    for insight in BCG_DEPT_INSIGHTS:
        doc.add_paragraph(insight, style="List Bullet")

    # ── Major-Level Analysis ──
    _add_heading(doc, "3. Major-Level Analysis (Enrollment-Based)", 2)
    _add_para(doc, (
        "The major-level view maps 48 individual majors using 2024 enrollment headcount "
        "as the market share proxy and 2022\u20132024 percentage change as the growth rate. "
        "Programs with fewer than 20 students in 2022 are flagged as 'small base' because their "
        "percentage changes can be misleading (e.g., Music: 12 to 27 = +125% from just 15 students)."
    ))
    _add_heading(doc, "Major-Level Quadrant Analysis", 3)
    quadrants = {
        "Stars (Large & Growing)": "Star",
        "Cash Cows (Large & Declining)": "Cash Cow",
        "Question Marks (Small & Growing)": "Question Mark",
        "Concerns (Small & Declining)": "Concern",
    }
    for label, q in quadrants.items():
        progs = BCG_DATA[BCG_DATA["Quadrant"] == q].sort_values("Enrollment_2024", ascending=False)
        _add_heading(doc, f"{label}: {len(progs)} majors", 3)
        for _, row in progs.head(10).iterrows():
            flag = " [small base]" if row["Small_Base"] else ""
            doc.add_paragraph(
                f"{row['Major']}: {int(row['Enrollment_2024'])} enrolled, "
                f"{row['Pct_Change']:+.1f}% change ({int(row['Abs_Change']):+d} students){flag}",
                style="List Bullet"
            )
        if len(progs) > 10:
            doc.add_paragraph(f"... and {len(progs) - 10} additional majors", style="List Bullet")

    _add_heading(doc, "4. Major-Level Key Findings", 2)
    for insight in BCG_INSIGHTS:
        doc.add_paragraph(insight, style="List Bullet")

    _add_heading(doc, "5. Portfolio Health Assessment", 2)
    counts = BCG_DATA["Quadrant"].value_counts()
    total = len(BCG_DATA)
    _add_para(doc, (
        f"FLC's academic portfolio of {total} majors shows the following distribution: "
        f"{counts.get('Star', 0)} Stars ({counts.get('Star', 0)*100//total}%), "
        f"{counts.get('Cash Cow', 0)} Cash Cows ({counts.get('Cash Cow', 0)*100//total}%), "
        f"{counts.get('Question Mark', 0)} Question Marks ({counts.get('Question Mark', 0)*100//total}%), "
        f"and {counts.get('Concern', 0)} Concerns ({counts.get('Concern', 0)*100//total}%). "
        f"Overall enrollment declined 3.1% from 2,899 to 2,810 between 2022 and 2024."
    ))
    _add_para(doc, (
        "The Concern quadrant contains many small programs. However, some (e.g., Philosophy, "
        "Native American & Indigenous Studies) contribute to FLC's liberal arts mission and "
        "general education requirements. NAIS is mission-critical and must not be evaluated "
        "solely on enrollment metrics. Strategic review should balance market data with mission "
        "contribution and interdisciplinary value."
    ))

    _add_heading(doc, "6. Investment Recommendations", 2)
    recommendations = [
        "Stars: Invest to sustain growth. Business Administration, Exercise Physiology, and "
        "Environmental Conservation & Mgmt lead the portfolio. Protect capacity and expand pathways.",
        "Cash Cows: Optimize for efficiency. Psychology and Biology/CMB maintain large enrollments "
        "but face declining trajectories. Focus on retention and right-sizing.",
        "Question Marks: Evaluate growth sustainability. Anthropology and Music show dramatic "
        "growth but from small bases. Validate whether trends are durable before major investment.",
        "Concerns: Conduct structured program review. Differentiate mission-critical programs "
        "(e.g., NAIS) from those that may benefit from consolidation or interdisciplinary pairing.",
        "Small-base caution: 12 programs had fewer than 20 students in 2022. Percentage changes "
        "for these programs should be interpreted with care.",
    ]
    for rec in recommendations:
        doc.add_paragraph(rec, style="List Bullet")

    _add_para(doc, "\nSource: Dataset_Majors.xlsx, FLC Institutional Data (2022\u20132024)", size=9)

    path = os.path.join(OUTPUT_DIR, "BCG_Executive_Summary.docx")
    doc.save(path)
    return path


def generate_bcg_pptx():
    prs = Presentation()
    _pptx_title_slide(prs, "BCG Growth-Share Matrix",
                       "Department & Major-Level Portfolio Analysis\nFort Lewis College")

    _pptx_content_slide(prs, "Methodology", [
        FRAMEWORK_DESCRIPTIONS["BCG"],
        "Department view: X = % of Total SCH, Y = 2-Year Change % (22 departments)",
        "Major view: X = 2024 Enrollment, Y = % Change 2022\u20132024 (48 majors)",
        "Bubble size (major view): Absolute enrollment change (students gained/lost)",
        "Small-base flag: Programs with < 20 students in 2022",
    ])

    # Department-level slides
    dept_counts = BCG_DEPT_DATA["Quadrant"].value_counts()
    _pptx_content_slide(prs, "Department-Level BCG (SCH-Based)", [
        f"Stars: {dept_counts.get('Star', 0)} departments (Business Admin, Psychology)",
        f"Cash Cows: {dept_counts.get('Cash Cow', 0)} departments (English, Math, Biology, HHP, Sociology, etc.)",
        f"Question Marks: {dept_counts.get('Question Mark', 0)} departments (Accounting, History)",
        f"Concerns: {dept_counts.get('Concern', 0)} departments (Political Science, Economics, Art & Design, etc.)",
        "Cash Cows generate bulk of institutional SCH but all show 2-year declines",
        "Source: BCG Presentation.pptx (FLC Internal)",
    ])

    # Major-level slides

    for label, q in [("Stars \u2014 Large & Growing", "Star"),
                      ("Cash Cows \u2014 Large & Declining", "Cash Cow"),
                      ("Question Marks \u2014 Small & Growing", "Question Mark"),
                      ("Concerns \u2014 Small & Declining", "Concern")]:
        progs = BCG_DATA[BCG_DATA["Quadrant"] == q].sort_values("Enrollment_2024", ascending=False)
        bullets = [f"{len(progs)} majors in this quadrant"]
        for _, row in progs.head(6).iterrows():
            flag = " *" if row["Small_Base"] else ""
            bullets.append(
                f"{row['Major']}: {int(row['Enrollment_2024'])} enrolled, "
                f"{row['Pct_Change']:+.1f}% ({int(row['Abs_Change']):+d}){flag}"
            )
        _pptx_content_slide(prs, label, bullets)

    counts = BCG_DATA["Quadrant"].value_counts()
    _pptx_content_slide(prs, "Portfolio Health", [
        f"{counts.get('Star', 0)} Stars, {counts.get('Cash Cow', 0)} Cash Cows, "
        f"{counts.get('Question Mark', 0)} Question Marks, {counts.get('Concern', 0)} Concerns",
        "Overall enrollment declined 3.1% (2,899 \u2192 2,810) from 2022 to 2024",
        "12 programs flagged as small base (< 20 students in 2022)",
        "Mission-critical programs (e.g., NAIS) must not be evaluated on enrollment alone",
    ])

    _pptx_content_slide(prs, "Investment Recommendations", [
        "Stars: Invest to sustain \u2014 Business Admin, Exercise Physiology, Env Conservation lead",
        "Cash Cows: Optimize efficiency \u2014 Psychology, Biology/CMB need retention focus",
        "Question Marks: Validate growth sustainability before major investment",
        "Concerns: Structured review; protect mission-critical programs",
        "Small-base caution: Interpret percentage changes for tiny programs with care",
    ])

    path = os.path.join(OUTPUT_DIR, "BCG_Slide_Deck.pptx")
    prs.save(path)
    return path


# ============================================================================
# SWOT MATRIX SLIDE
# ============================================================================

def _swot_slide_title(slide, title_text):
    """Add FLC-branded title bar to a SWOT slide."""
    title_box = slide.shapes.add_textbox(
        PptxInches(0.4), PptxInches(0.25), PptxInches(12.5), PptxInches(0.6),
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    run.font.size = PptxPt(28)
    run.font.color.rgb = PptxRGB(0x00, 0x30, 0x57)
    run.font.bold = True
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    run2 = p2.add_run()
    run2.text = (
        f"Fort Lewis College | {INSTITUTION['location']} | "
        "Cross-Framework Strategic Synthesis"
    )
    run2.font.size = PptxPt(11)
    run2.font.color.rgb = PptxRGB(0x66, 0x66, 0x66)


def _swot_quadrant(slide, label, color, left, top, box_width, box_height):
    """Render one SWOT quadrant on a slide."""
    data = SWOT_DATA[label]

    # Colored header bar
    header = slide.shapes.add_textbox(
        PptxInches(left), PptxInches(top), box_width, PptxInches(0.38),
    )
    header.fill.solid()
    header.fill.fore_color.rgb = color
    htf = header.text_frame
    htf.word_wrap = True
    htf.margin_top = PptxPt(3)
    htf.margin_bottom = PptxPt(3)
    htf.margin_left = PptxPt(8)
    hp = htf.paragraphs[0]
    hp.alignment = PP_ALIGN.LEFT
    hr = hp.add_run()
    hr.text = f"{label}  ({len(data['items'])} items)"
    hr.font.size = PptxPt(13)
    hr.font.bold = True
    hr.font.color.rgb = PptxRGB(0xFF, 0xFF, 0xFF)

    # Content body
    body = slide.shapes.add_textbox(
        PptxInches(left), PptxInches(top + 0.38),
        box_width, PptxInches(box_height - 0.38),
    )
    body.fill.solid()
    body.fill.fore_color.rgb = PptxRGB(0xFA, 0xFA, 0xFA)
    body.line.color.rgb = PptxRGB(0xE0, 0xE0, 0xE0)
    body.line.width = PptxPt(0.5)
    btf = body.text_frame
    btf.word_wrap = True
    btf.margin_top = PptxPt(6)
    btf.margin_left = PptxPt(8)
    btf.margin_right = PptxPt(6)
    btf.vertical_anchor = MSO_ANCHOR.TOP

    for idx, item in enumerate(data["items"]):
        # Title line (bold)
        if idx == 0:
            p_title = btf.paragraphs[0]
        else:
            p_title = btf.add_paragraph()
        p_title.space_before = PptxPt(4) if idx > 0 else PptxPt(0)
        p_title.space_after = PptxPt(1)
        r_bullet = p_title.add_run()
        r_bullet.text = f"\u2022 {item['title']}"
        r_bullet.font.size = PptxPt(9)
        r_bullet.font.bold = True
        r_bullet.font.color.rgb = PptxRGB(0x00, 0x30, 0x57)

        # Detail line
        p_detail = btf.add_paragraph()
        p_detail.space_before = PptxPt(0)
        p_detail.space_after = PptxPt(1)
        r_detail = p_detail.add_run()
        r_detail.text = f"   {item['detail']}"
        r_detail.font.size = PptxPt(7.5)
        r_detail.font.color.rgb = PptxRGB(0x33, 0x33, 0x33)

        # Source line (italic)
        p_src = btf.add_paragraph()
        p_src.space_before = PptxPt(0)
        p_src.space_after = PptxPt(2)
        r_src = p_src.add_run()
        r_src.text = f"   Source: {item['source']}"
        r_src.font.size = PptxPt(6.5)
        r_src.font.italic = True
        r_src.font.color.rgb = PptxRGB(0x99, 0x99, 0x99)


def generate_swot_pptx():
    """Generate two landscape SWOT slides: Internal Factors and External Factors."""
    prs = Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)

    box_width = PptxInches(6.1)
    col_left = 0.4
    col_right = 6.9
    content_top = 1.15
    content_height = 5.95  # full height below title area

    # --- Slide 1: Internal Factors (Strengths + Weaknesses) ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    _swot_slide_title(slide1, "SWOT Analysis \u2014 Internal Factors")
    _swot_quadrant(slide1, "Strengths", PptxRGB(0x2E, 0xCC, 0x71),
                   col_left, content_top, box_width, content_height)
    _swot_quadrant(slide1, "Weaknesses", PptxRGB(0xE7, 0x4C, 0x3C),
                   col_right, content_top, box_width, content_height)

    # --- Slide 2: External Factors (Opportunities + Threats) ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    _swot_slide_title(slide2, "SWOT Analysis \u2014 External Factors")
    _swot_quadrant(slide2, "Opportunities", PptxRGB(0x34, 0x98, 0xDB),
                   col_left, content_top, box_width, content_height)
    _swot_quadrant(slide2, "Threats", PptxRGB(0xE6, 0x7E, 0x22),
                   col_right, content_top, box_width, content_height)

    path = os.path.join(OUTPUT_DIR, "SWOT_Matrix.pptx")
    prs.save(path)
    return path


# ============================================================================
# PROJECT-LEVEL SYNTHESIS DOCUMENTS
# ============================================================================

def _bcg_quadrant_counts():
    """Return dict of quadrant label -> count."""
    counts = BCG_DATA["Quadrant"].value_counts().to_dict()
    return {
        "Stars": counts.get("Star", 0),
        "Cash Cows": counts.get("Cash Cow", 0),
        "Question Marks": counts.get("Question Mark", 0),
        "Concerns": counts.get("Concern", 0),
    }


def _gray_rec_counts():
    """Return dict of recommendation -> count."""
    return GRAY_ASSOCIATES_DATA["GA_Recommendation"].value_counts().to_dict()


def _write_exec_summary_content(doc):
    """Shared executive summary narrative used by both the standalone and the final report."""
    bcg_counts = _bcg_quadrant_counts()
    gray_counts = _gray_rec_counts()

    # Project Overview
    _add_heading(doc, "Project Overview", 2)
    _add_para(doc, (
        f"The Fort Lewis College Portfolio Optimization Project is a comprehensive strategic analysis "
        f"of FLC's academic portfolio, designed to inform the 2025\u20132030 planning cycle. "
        f"FLC enrolls {INSTITUTION['total_enrollment_f25']:,} students in Durango, Colorado, with a "
        f"{INSTITUTION['retention_rate_f24']}% first-time full-time retention rate. As a small, public, "
        f"rural institution with a Native American\u2013serving mission, FLC faces unique competitive "
        f"pressures alongside distinctive strengths."
    ))
    _add_para(doc, (
        "The project applies a three-phase methodology: Phase 1 (Environmental Scanning) uses four "
        "frameworks\u2014PESTLE, Porter's Five Forces, BCG Growth-Share Matrix, and Gray Associates "
        "Portfolio Analysis\u2014to assess external conditions and internal portfolio health. Phase 2 "
        "(Strategic Synthesis) consolidates findings into a SWOT analysis. Phase 3 (Strategic Direction) "
        "applies the Zone to Win framework to produce three actionable scenarios with a strategic roadmap."
    ))

    # Environmental Scanning
    _add_heading(doc, "Environmental Scanning Findings", 2)
    # PESTLE
    _add_heading(doc, "PESTLE Analysis", 3)
    top_pestle = sorted(PESTLE_DATA.items(), key=lambda x: x[1]["impact_score"], reverse=True)[:3]
    pestle_summary = ", ".join([f"{cat} ({d['impact_score']}/5)" for cat, d in top_pestle])
    _add_para(doc, (
        f"The PESTLE analysis assessed six external dimensions. Highest-impact categories: {pestle_summary}. "
        f"Federal DEI policy disruption threatens TRIO programs (120+ terminated nationally) and could "
        f"misclassify FLC's statutory NATW mission. State funding is projected down for FY2026. "
        f"The Durango housing crisis constrains faculty/staff recruitment. Online market is saturated\u2014"
        f"FLC's only defensible online strategy is the Indigenous niche (NATW tuition waiver moat)."
    ))

    # Porter's
    _add_heading(doc, "Porter's Five Forces", 3)
    avg_score = sum(d["score"] for d in PORTERS_DATA.values()) / len(PORTERS_DATA)
    _add_para(doc, (
        f"Overall competitive intensity is HIGH ({avg_score:.1f}/5 average across five forces). "
        f"Competitive rivalry ({PORTERS_DATA['Competitive Rivalry']['score']}/5) and student bargaining "
        f"power ({PORTERS_DATA['Bargaining Power of Students']['score']}/5) are the strongest pressures. "
        f"FLC's defensive moats include its Native American tuition waiver, Durango's outdoor lifestyle, "
        f"and an intimate small-college experience."
    ))

    # BCG
    _add_heading(doc, "BCG Growth-Share Matrix", 3)
    _total_bcg = sum(bcg_counts.values())
    _dept_counts = BCG_DEPT_DATA["Quadrant"].value_counts()
    _add_para(doc, (
        f"The BCG analysis examines FLC at two levels. At the department level (SCH-based), 22 departments "
        f"show {_dept_counts.get('Star', 0)} Stars, {_dept_counts.get('Cash Cow', 0)} Cash Cows, "
        f"and {_dept_counts.get('Concern', 0)} Concerns \u2014 Cash Cows generate bulk revenue but all "
        f"face declining enrollment. At the major level (enrollment-based), {_total_bcg} majors show "
        f"{bcg_counts['Stars']} Stars, {bcg_counts['Cash Cows']} Cash Cows, "
        f"{bcg_counts['Question Marks']} Question Marks, and {bcg_counts['Concerns']} Concerns. "
        f"Programs with fewer than 20 students in 2022 are flagged as small-base."
    ))

    # Gray
    _add_heading(doc, "Gray Associates Portfolio Analysis", 3)
    rec_str = ", ".join([f"{count} {rec}" for rec, count in sorted(gray_counts.items(),
                         key=lambda x: ["Grow", "Sustain", "Transform", "Evaluate", "Sunset Review"].index(x[0]))])
    _add_para(doc, (
        f"The Gray Associates framework evaluated 23 degree programs on Market Score and Economics Score "
        f"(note: applied with proxy data; scores are directional, not definitive): "
        f"{rec_str}. Programs recommended to Grow (Engineering, Health Sciences, Business, Psychology) "
        f"align with strong employment demand and student interest."
    ))

    # SWOT
    _add_heading(doc, "Strategic Synthesis (SWOT)", 2)
    swot_counts = {q: len(d["items"]) for q, d in SWOT_DATA.items()}
    _add_para(doc, (
        f"The SWOT analysis synthesized all Phase 1 findings: {swot_counts['Strengths']} Strengths, "
        f"{swot_counts['Weaknesses']} Weaknesses, {swot_counts['Opportunities']} Opportunities, "
        f"and {swot_counts['Threats']} Threats. Top strengths: statutory NATW mission (CRS 23-52-105) "
        f"and 15 Star-quadrant majors. Critical tension reconciled: Indigenous education represents "
        f"FLC's best growth opportunity but must be framed through statutory obligations (CRS 23-52-105), not DEI language, "
        f"given current federal policy climate."
    ))

    # Zone to Win + Scenarios
    _add_heading(doc, "Strategic Direction: Zone to Win", 2)
    _add_para(doc, (
        "The Zone to Win framework (Geoffrey Moore) is the culminating analytical step, translating "
        "all prior findings into four strategic investment zones\u2014Performance, Productivity, "
        "Incubation, and Transformation\u2014across three scenarios."
    ))
    for zone, zd in ZONE_TO_WIN_DATA.items():
        _add_para(doc, f"{zone}: {zd['description']}", bold=True, size=10)

    _add_heading(doc, "Three Strategic Scenarios", 3)
    for name, sd in SCENARIOS.items():
        _add_para(doc, f"{name} Scenario", bold=True, size=11)
        _add_para(doc, (
            f"{sd['description']} Enrollment target: {sd['enrollment_target']:,}. "
            f"Retention target: {sd['retention_target']}%. "
            f"Zone allocation: Performance {sd['zone_allocation']['Performance']}%, "
            f"Productivity {sd['zone_allocation']['Productivity']}%, "
            f"Incubation {sd['zone_allocation']['Incubation']}%, "
            f"Transformation {sd['zone_allocation']['Transformation']}%."
        ))

    # Recommendations
    _add_heading(doc, "Recommendations", 2)
    recommendations = [
        "Invest in Star-quadrant majors (Business Admin, Exercise Physiology, Psychology, Engineering) with capacity and pipeline growth",
        "Initiate structured review of 17 Concern-quadrant majors, balancing enrollment data with mission alignment (protect NAIS regardless)",
        "Prioritize retention improvement as most cost-effective enrollment strategy (Compass expansion, early-alert, advising redesign)",
        "Pursue Indigenous education as a defensible national niche\u2014frame through statutory/sovereign obligations, not DEI; invest in marketing for NATW online recruitment",
        "Strengthen dual enrollment and transfer pipeline as a near-term enrollment stabilizer (3+ high school partnerships)",
        "Address Durango housing crisis impact on faculty/staff recruitment through institutional housing partnerships",
        "Navigate faculty governance deliberately\u2014accept 12-18 month timelines for program restructuring; build evidence base with BCG and Gray data",
    ]
    for rec in recommendations:
        doc.add_paragraph(rec, style="List Bullet")


def generate_exec_summary_docx():
    """Generate 3-4 page narrative executive summary pulling from all frameworks."""
    doc = Document()
    _doc_header(doc, "Executive Summary",
                "FLC Portfolio Optimization Project \u2014 Strategic Analysis & Recommendations")

    _write_exec_summary_content(doc)

    _add_para(doc, "\nSource: FLC Portfolio Optimization Project \u2014 PESTLE, Porter's, BCG, "
              "Gray Associates, SWOT, and Zone to Win analyses", size=9)

    path = os.path.join(OUTPUT_DIR, "Executive_Summary_Narrative.docx")
    doc.save(path)
    return path


def generate_exec_summary_pptx():
    """Generate 5-8 slide executive summary deck with high points from all frameworks."""
    prs = Presentation()
    bcg_counts = _bcg_quadrant_counts()
    gray_counts = _gray_rec_counts()
    _dept_counts = BCG_DEPT_DATA["Quadrant"].value_counts()

    _pptx_title_slide(prs, "FLC Portfolio Optimization Project",
                       "Executive Summary\nFort Lewis College Academic Affairs")

    # Slide 2: Project At-A-Glance
    _pptx_content_slide(prs, "Project At-A-Glance", [
        f"Enrollment: {INSTITUTION['total_enrollment_f25']:,} students (Fall 2025)",
        f"Retention: {INSTITUTION['retention_rate_f24']}% FTFT (Fall 2024)",
        "Graduate Students: 160 (1 existing program, Fall 2025)",
        "BCG: 22 departments (SCH-based) + 48 majors (enrollment-based); Gray: 23 programs",
        "3-Phase Methodology: Environmental Scanning \u2192 SWOT Synthesis \u2192 Zone to Win Scenarios",
    ])

    # Slide 3: Phase 1 Key Findings
    top_pestle = sorted(PESTLE_DATA.items(), key=lambda x: x[1]["impact_score"], reverse=True)[:3]
    pestle_line = ", ".join([f"{cat} ({d['impact_score']}/5)" for cat, d in top_pestle])
    avg_porters = sum(d["score"] for d in PORTERS_DATA.values()) / len(PORTERS_DATA)
    _pptx_content_slide(prs, "Phase 1: Environmental Scanning", [
        f"PESTLE: Top impacts \u2014 {pestle_line}",
        f"Porter's: Overall competitive intensity {avg_porters:.1f}/5 (HIGH)",
        f"BCG: Depts \u2014 {_dept_counts.get('Star', 0)} Stars, {_dept_counts.get('Cash Cow', 0)} Cash Cows, {_dept_counts.get('Concern', 0)} Concerns; "
        f"Majors \u2014 {bcg_counts['Stars']} Stars, {bcg_counts['Concerns']} Concerns",
        f"Gray: {gray_counts.get('Grow', 0)} Grow, {gray_counts.get('Sustain', 0)} Sustain, "
        f"{gray_counts.get('Transform', 0)} Transform, {gray_counts.get('Evaluate', 0)} Evaluate, "
        f"{gray_counts.get('Sunset Review', 0)} Sunset Review",
    ])

    # Slide 4: SWOT Highlights
    swot_bullets = []
    for quadrant in ["Strengths", "Weaknesses", "Opportunities", "Threats"]:
        items = SWOT_DATA[quadrant]["items"][:2]
        titles = " | ".join([it["title"] for it in items])
        swot_bullets.append(f"{quadrant}: {titles}")
    _pptx_content_slide(prs, "Phase 2: SWOT Synthesis", swot_bullets)

    # Slide 5: Zone to Win
    zone_bullets = []
    for zone, zd in ZONE_TO_WIN_DATA.items():
        zone_bullets.append(f"{zone}: {len(zd['programs'])} initiatives \u2014 {zd['programs'][0]['name']}, {zd['programs'][1]['name']}")
    _pptx_content_slide(prs, "Phase 3: Zone to Win Framework", zone_bullets)

    # Slide 6: Scenarios
    scenario_bullets = []
    for name, sd in SCENARIOS.items():
        scenario_bullets.append(
            f"{name}: Enrollment {sd['enrollment_target']:,}, Retention {sd['retention_target']}%, "
            f"{sd['new_programs']} new programs"
        )
    _pptx_content_slide(prs, "Three Strategic Scenarios", scenario_bullets)

    # Slide 7: Recommendations
    _pptx_content_slide(prs, "Key Recommendations", [
        "Invest in Star majors: Business Admin, Exercise Physiology, Psychology, Engineering",
        "Initiate structured review of 17 Concern-quadrant majors (protect NAIS regardless)",
        "Prioritize retention as most cost-effective enrollment strategy (Compass, advising, early-alert)",
        "Pursue Indigenous online niche (statutorily grounded, NATW moat) \u2014 not generic online degrees",
        "Strengthen dual enrollment (3+ schools) and transfer pipelines",
        "Address Durango housing barrier to faculty/staff recruitment",
    ])

    path = os.path.join(OUTPUT_DIR, "Executive_Summary_Deck.pptx")
    prs.save(path)
    return path


def generate_final_report_docx():
    """Generate ~8-10 page final report with Executive Summary, Analysis, Risk, and Summary sections."""
    doc = Document()
    _doc_header(doc, "Final Report",
                "FLC Portfolio Optimization Project \u2014 Comprehensive Strategic Analysis")

    # ── SECTION 1: EXECUTIVE SUMMARY (2 pages) ──
    _add_heading(doc, "Section 1: Executive Summary", 1)
    _write_exec_summary_content(doc)

    # ── SECTION 2: ANALYSIS (4-6 pages) ──
    _add_heading(doc, "Section 2: Detailed Analysis", 1)
    _add_para(doc, (
        "This section presents 3\u20135 major takeaways from each Phase 1 framework and the "
        "Phase 2 SWOT synthesis. Zone to Win is excluded here as it is the culminating "
        "framework covered in the Executive Summary."
    ))

    # PESTLE
    _add_heading(doc, "PESTLE Analysis", 2)
    sorted_pestle = sorted(PESTLE_DATA.items(), key=lambda x: x[1]["impact_score"], reverse=True)
    for cat, d in sorted_pestle:
        _add_heading(doc, f"{cat} \u2014 Impact: {d['impact']} ({d['impact_score']}/5), Trend: {d['trend']}", 3)
        for factor in d["factors"][:3]:
            doc.add_paragraph(factor, style="List Bullet")
        if d["opportunities"]:
            doc.add_paragraph(f"Opportunity: {d['opportunities'][0]}", style="List Bullet 2")

    # Porter's
    _add_heading(doc, "Porter's Five Forces", 2)
    for force, d in PORTERS_DATA.items():
        _add_heading(doc, f"{force}: {d['rating']} ({d['score']}/5)", 3)
        _add_para(doc, d["description"])
        for ind in d["indicators"][:3]:
            doc.add_paragraph(
                f"{ind['name']}: {ind['value']} (Trend: {ind['trend']})",
                style="List Bullet"
            )
    _add_para(doc, "Key Takeaways:", bold=True, size=10)
    for insight in PORTERS_INSIGHTS:
        doc.add_paragraph(insight, style="List Bullet")

    # BCG
    _add_heading(doc, "BCG Growth-Share Matrix", 2)
    _add_para(doc, (
        "The BCG analysis examines FLC's portfolio at two levels: 22 departments (SCH share as market "
        "share proxy) and 48 individual majors (2024 enrollment headcount vs. 2022\u20132024 % change)."
    ))

    # Department-level
    _add_heading(doc, "Department-Level Analysis (SCH-Based)", 3)
    dept_quadrant_labels = {"Star": "Stars", "Cash Cow": "Cash Cows",
                            "Question Mark": "Question Marks", "Concern": "Concerns"}
    for q_key, q_label in dept_quadrant_labels.items():
        depts = BCG_DEPT_DATA[BCG_DEPT_DATA["Quadrant"] == q_key]
        _add_para(doc, f"{q_label} ({len(depts)} departments):", bold=True, size=10)
        for _, row in depts.iterrows():
            doc.add_paragraph(
                f"{row['Department']}: {row['SCH_Pct']}% of SCH, {row['Two_Year_Change']:+.1f}% 2-year change",
                style="List Bullet"
            )

    # Major-level
    _add_heading(doc, "Major-Level Analysis (Enrollment-Based)", 3)
    bcg_counts = _bcg_quadrant_counts()
    quadrant_labels = {"Star": "Stars", "Cash Cow": "Cash Cows",
                       "Question Mark": "Question Marks", "Concern": "Concerns"}
    for q_key, q_label in quadrant_labels.items():
        majors = BCG_DATA[BCG_DATA["Quadrant"] == q_key]
        _add_para(doc, f"{q_label} ({len(majors)} majors):", bold=True, size=10)
        for _, row in majors.head(8).iterrows():
            flag = " [Small Base]" if row["Small_Base"] else ""
            doc.add_paragraph(
                f"{row['Major']}: {row['Enrollment_2024']} students, {row['Pct_Change']:+.1f}% change{flag}",
                style="List Bullet"
            )
        if len(majors) > 8:
            doc.add_paragraph(f"... and {len(majors) - 8} additional majors", style="List Bullet")
    _add_para(doc, "Key Takeaways:", bold=True, size=10)
    for insight in BCG_INSIGHTS:
        doc.add_paragraph(insight, style="List Bullet")

    # Gray Associates
    _add_heading(doc, "Gray Associates Portfolio Analysis", 2)
    gray_counts = _gray_rec_counts()
    for rec in ["Grow", "Sustain", "Transform", "Evaluate", "Sunset Review"]:
        progs = GRAY_ASSOCIATES_DATA[GRAY_ASSOCIATES_DATA["GA_Recommendation"] == rec]
        if not progs.empty:
            _add_heading(doc, f"{rec} ({len(progs)} programs)", 3)
            for _, p in progs.head(5).iterrows():
                doc.add_paragraph(
                    f"{p['Program']}: Market {p['Market_Score']}, Economics {p['Economics_Score']}, "
                    f"Enrollment {p['Enrollment']}",
                    style="List Bullet"
                )
    _add_para(doc, "Key Takeaways:", bold=True, size=10)
    for insight in GA_INSIGHTS:
        doc.add_paragraph(insight, style="List Bullet")

    # SWOT
    _add_heading(doc, "SWOT Synthesis", 2)
    for quadrant in ["Strengths", "Weaknesses", "Opportunities", "Threats"]:
        items = SWOT_DATA[quadrant]["items"]
        _add_heading(doc, f"{quadrant} ({len(items)} items)", 3)
        for item in items[:4]:
            doc.add_paragraph(f"{item['title']}: {item['detail']}", style="List Bullet")

    # ── SECTION 3: RISK AND IMPLEMENTATION (1/2 page) ──
    _add_heading(doc, "Section 3: Risk & Implementation", 1)
    _add_para(doc, (
        f"The strategic roadmap includes {len(ROADMAP_MILESTONES)} milestones tracked across "
        f"all three phases, with {len(ROADMAP_KPIS)} key performance indicators measuring progress. "
        f"Below are the highest-priority risks and mitigation strategies."
    ))
    # Risk table
    table = doc.add_table(rows=1, cols=4)
    table.style = "Light Grid Accent 1"
    hdr = table.rows[0].cells
    hdr[0].text = "Risk"
    hdr[1].text = "Probability"
    hdr[2].text = "Impact"
    hdr[3].text = "Mitigation Strategy"
    for _, row in RISK_MITIGATION.head(6).iterrows():
        cells = table.add_row().cells
        cells[0].text = row["Risk"]
        cells[1].text = row["Probability"]
        cells[2].text = row["Impact"]
        cells[3].text = row["Mitigation_Strategy"]

    # ── SECTION 4: SUMMARY (1 page) ──
    _add_heading(doc, "Section 4: Summary", 1)
    _add_para(doc, (
        "Fort Lewis College operates in an increasingly competitive and financially constrained "
        "higher education environment. Declining college-going rates, intensifying online competition, "
        "state funding volatility, and political pressure on diversity programs create significant "
        "headwinds. At the same time, FLC possesses distinctive strengths that few competitors can "
        "replicate: a federally rooted Native American tuition waiver, an unmatched outdoor recreation "
        "setting in Durango, and an intimate small-college experience."
    ))
    _add_para(doc, (
        "The analysis reveals a portfolio in need of strategic rebalancing. Of 48 majors analyzed, "
        "17 fall in the BCG Concern quadrant while 15 are Stars, indicating both risk and opportunity. "
        "The Gray Associates framework identifies programs for growth investment, sustaining, and those "
        "requiring transformation, evaluation, or sunset review."
    ))
    _add_para(doc, (
        "The Zone to Win framework provides the decision architecture for action. Three qualitatively "
        "different scenarios (Incremental, Moderate-Adaptive, Disruptive) define distinct strategic bets "
        "rather than simply varying the degree of investment. Consistent priorities across all scenarios: "
        "invest in proven Star programs, improve retention through advising and early-alert systems, "
        "differentiate through Indigenous education (statutorily grounded, not DEI), and leverage FLC's unique "
        "place-based and experiential learning assets."
    ))
    _add_para(doc, (
        "Success requires navigating faculty governance deliberately (12-18 month timelines), addressing "
        "the Durango housing barrier to recruitment, investing marketing resources in the Indigenous "
        "student niche (the only defensible online strategy), and maintaining honest assessment of what "
        "is achievable within FLC's constrained budget."
    ))

    _add_para(doc, "\nSource: FLC Portfolio Optimization Project \u2014 All framework analyses, "
              "institutional data, and strategic planning documents", size=9)

    path = os.path.join(OUTPUT_DIR, "Final_Report.docx")
    doc.save(path)
    return path


def generate_final_report_pptx():
    """Generate max 20-slide final presentation leveraging Final Report content."""
    prs = Presentation()
    bcg_counts = _bcg_quadrant_counts()
    gray_counts = _gray_rec_counts()
    _dept_counts = BCG_DEPT_DATA["Quadrant"].value_counts()

    # Slide 1: Title
    _pptx_title_slide(prs, "FLC Portfolio Optimization",
                       "Final Presentation\nFort Lewis College Academic Affairs")

    # Slide 2: Agenda
    _pptx_content_slide(prs, "Agenda", [
        "Project Overview & Methodology",
        "Phase 1: Environmental Scanning (PESTLE, Porter's, BCG, Gray Associates)",
        "Phase 2: Strategic Synthesis (SWOT)",
        "Phase 3: Strategic Direction (Zone to Win & Scenarios)",
        "Risk & Implementation",
        "Recommendations & Next Steps",
    ])

    # Slide 3: Project Overview
    _pptx_content_slide(prs, "Project Overview", [
        f"Fort Lewis College \u2014 {INSTITUTION['location']} | {INSTITUTION['type']}",
        f"Enrollment: {INSTITUTION['total_enrollment_f25']:,} | Retention: {INSTITUTION['retention_rate_f24']}%",
        "Native American\u2013serving mission with federally rooted tuition waiver",
        "3-Phase methodology: Scanning \u2192 Synthesis \u2192 Direction",
        "7 frameworks; BCG: 22 departments (SCH) + 48 majors (enrollment); Gray: 23 programs",
    ])

    # Slide 4: PESTLE Key Findings
    sorted_pestle = sorted(PESTLE_DATA.items(), key=lambda x: x[1]["impact_score"], reverse=True)
    pestle_bullets = []
    for cat, d in sorted_pestle[:4]:
        pestle_bullets.append(f"{cat}: {d['impact']} ({d['impact_score']}/5) \u2014 {d['factors'][0]}")
    _pptx_content_slide(prs, "PESTLE: Key Findings", pestle_bullets)

    # Slide 5: PESTLE Implications
    pestle_opps = []
    for cat, d in sorted_pestle[:4]:
        if d["opportunities"]:
            pestle_opps.append(f"{cat}: {d['opportunities'][0]}")
    _pptx_content_slide(prs, "PESTLE: Opportunities", pestle_opps)

    # Slide 6: Porter's Overview
    porters_bullets = []
    for force, d in PORTERS_DATA.items():
        porters_bullets.append(f"{force}: {d['rating']} ({d['score']}/5)")
    avg_porters = sum(d["score"] for d in PORTERS_DATA.values()) / len(PORTERS_DATA)
    porters_bullets.append(f"Overall intensity: {avg_porters:.1f}/5 (HIGH)")
    _pptx_content_slide(prs, "Porter's Five Forces", porters_bullets)

    # Slide 7: Porter's Competitive Position
    _pptx_content_slide(prs, "Porter's: Competitive Position", PORTERS_INSIGHTS[:5])

    # Slide 8: BCG Department Summary
    _pptx_content_slide(prs, "BCG: Department-Level (SCH-Based)", [
        f"Stars: {_dept_counts.get('Star', 0)} departments (Business Admin, Psychology)",
        f"Cash Cows: {_dept_counts.get('Cash Cow', 0)} departments (English, Math, Biology, HHP, etc.)",
        f"Question Marks: {_dept_counts.get('Question Mark', 0)} departments (Accounting, History)",
        f"Concerns: {_dept_counts.get('Concern', 0)} departments (Poli Sci, Economics, Art & Design, etc.)",
        "Cash Cows generate bulk of institutional SCH but all show 2-year enrollment declines",
    ])

    # Slide 9: BCG Major Summary
    _total_bcg = sum(bcg_counts.values())
    _pptx_content_slide(prs, "BCG: Major-Level (Enrollment-Based)", [
        f"Stars: {bcg_counts['Stars']} majors (above-median enrollment, positive growth)",
        f"Cash Cows: {bcg_counts['Cash Cows']} majors (above-median enrollment, negative/flat growth)",
        f"Question Marks: {bcg_counts['Question Marks']} majors (below-median enrollment, positive growth)",
        f"Concerns: {bcg_counts['Concerns']} majors (below-median enrollment, negative/flat growth)",
        f"Portfolio note: {bcg_counts['Concerns']*100//_total_bcg}% in Concern quadrant; small-base programs flagged",
    ])

    # Slide 10: BCG Insights
    _pptx_content_slide(prs, "BCG: Key Takeaways", BCG_INSIGHTS[:5])

    # Slide 10: Gray Overview
    gray_bullets = []
    for rec in ["Grow", "Sustain", "Transform", "Evaluate", "Sunset Review"]:
        count = gray_counts.get(rec, 0)
        progs = GRAY_ASSOCIATES_DATA[GRAY_ASSOCIATES_DATA["GA_Recommendation"] == rec]
        top_names = ", ".join(progs["Program"].head(3).tolist())
        gray_bullets.append(f"{rec} ({count}): {top_names}")
    _pptx_content_slide(prs, "Gray Associates: Program Classifications", gray_bullets)

    # Slide 11: Gray Insights
    _pptx_content_slide(prs, "Gray Associates: Key Takeaways", GA_INSIGHTS[:5])

    # Slide 12: SWOT Summary
    swot_bullets = []
    for quadrant in ["Strengths", "Weaknesses", "Opportunities", "Threats"]:
        items = SWOT_DATA[quadrant]["items"][:2]
        titles = ", ".join([it["title"] for it in items])
        swot_bullets.append(f"{quadrant} ({len(SWOT_DATA[quadrant]['items'])}): {titles}")
    _pptx_content_slide(prs, "SWOT Synthesis", swot_bullets)

    # Slide 13: SWOT Strategic Themes
    swot_themes = []
    for quadrant in ["Strengths", "Weaknesses", "Opportunities", "Threats"]:
        top = SWOT_DATA[quadrant]["items"][0]
        swot_themes.append(f"{quadrant}: {top['title']} \u2014 {top['detail'][:80]}...")
    _pptx_content_slide(prs, "SWOT: Strategic Themes", swot_themes)

    # Slide 14: Zone to Win Framework
    zone_bullets = []
    for zone, zd in ZONE_TO_WIN_DATA.items():
        zone_bullets.append(f"{zone} ({len(zd['programs'])} initiatives): {zd['description'][:80]}...")
    _pptx_content_slide(prs, "Zone to Win Framework", zone_bullets)

    # Slide 15: Zone to Win Programs
    zone_prog_bullets = []
    for zone, zd in ZONE_TO_WIN_DATA.items():
        top = zd["programs"][:2]
        names = ", ".join([p["name"] for p in top])
        zone_prog_bullets.append(f"{zone}: {names}")
    _pptx_content_slide(prs, "Zone to Win: Key Programs", zone_prog_bullets)

    # Slide 16: Strategic Scenarios
    scenario_bullets = []
    for name, sd in SCENARIOS.items():
        scenario_bullets.append(
            f"{name}: Enrollment {sd['enrollment_target']:,}, Retention {sd['retention_target']}%, "
            f"{sd['new_programs']} new programs"
        )
        scenario_bullets.append(f"  \u2192 {sd['description'][:80]}")
    _pptx_content_slide(prs, "Three Strategic Scenarios", scenario_bullets)

    # Slide 17: Risk & Implementation
    risk_bullets = []
    for _, row in RISK_MITIGATION.head(5).iterrows():
        risk_bullets.append(f"{row['Risk']} ({row['Probability']}/{row['Impact']}): {row['Mitigation_Strategy'][:60]}...")
    _pptx_content_slide(prs, "Risk & Implementation", risk_bullets)

    # Slide 18: Recommendations
    _pptx_content_slide(prs, "Key Recommendations", [
        "Invest in Star majors (Business Admin, Exercise Physiology, Psychology, Engineering)",
        "Initiate structured review of 17 Concern-quadrant majors (protect NAIS regardless)",
        "Prioritize retention as most cost-effective enrollment strategy",
        "Pursue Indigenous online niche (statutorily grounded, NATW moat) \u2014 not generic degrees",
        "Strengthen dual enrollment (3+ schools) and transfer pipelines",
        "Address Durango housing barrier; navigate faculty governance with 12-18 month timelines",
    ])

    path = os.path.join(OUTPUT_DIR, "Final_Presentation.pptx")
    prs.save(path)
    return path


# ============================================================================
# PUBLIC API
# ============================================================================

GENERATORS = {
    "pestle_docx": generate_pestle_docx,
    "pestle_pptx": generate_pestle_pptx,
    "porters_docx": generate_porters_docx,
    "porters_pptx": generate_porters_pptx,
    "gray_docx": generate_gray_docx,
    "gray_pptx": generate_gray_pptx,
    "bcg_docx": generate_bcg_docx,
    "bcg_pptx": generate_bcg_pptx,
    "swot_pptx": generate_swot_pptx,
    "exec_summary_docx": generate_exec_summary_docx,
    "exec_summary_pptx": generate_exec_summary_pptx,
    "final_report_docx": generate_final_report_docx,
    "final_report_pptx": generate_final_report_pptx,
}


def generate_all():
    """Generate all documents. Returns dict of {name: filepath}."""
    results = {}
    for name, fn in GENERATORS.items():
        results[name] = fn()
        print(f"  Generated: {results[name]}")
    return results


if __name__ == "__main__":
    print("\nGenerating all executive summaries and slide decks...\n")
    generate_all()
    print(f"\nAll documents saved to: {OUTPUT_DIR}\n")
